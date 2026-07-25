"""Content extraction engine — dispatches to type-specific extractors.

Supports: PDF, DOCX, PPTX, XLSX, TXT, MD, HTML, CSV, images (metadata only).
"""

from pathlib import Path
from typing import Any
from loguru import logger

from ..models.file_item import FileItem, FileStatus


# ── Public API ──────────────────────────────────────────────────────────────

def extract(file_item: FileItem, max_chars: int = 4000) -> None:
    """Extract text content and metadata from a file.

    Updates file_item in place with extracted data and status.

    Args:
        file_item: The FileItem to process.
        max_chars: Maximum characters to extract (for AI API limits).
    """
    ext = file_item.original_ext
    path = file_item.original_path

    if not path.exists():
        file_item.status = FileStatus.FAILED
        file_item.error_message = "文件不存在"
        return

    file_item.status = FileStatus.EXTRACTING

    try:
        if ext == ".pdf":
            text, meta = _extract_pdf(path)
        elif ext == ".docx":
            text, meta = _extract_docx(path)
        elif ext == ".pptx":
            text, meta = _extract_pptx(path)
        elif ext in (".xlsx", ".xls"):
            text, meta = _extract_xlsx(path)
        elif ext in (".txt", ".md", ".csv", ".html", ".htm", ".rtf"):
            text, meta = _extract_text(path)
        elif ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".avif"):
            text, meta = _extract_image(path)
        elif ext in (".mp4", ".mov", ".mkv", ".avi", ".webm", ".wmv", ".flv"):
            text, meta = _extract_video(path)
        else:
            file_item.status = FileStatus.FAILED
            file_item.error_message = f"不支持的文件类型: {ext}"
            return

        # Truncate to max_chars for AI API
        file_item.extracted_text = text[:max_chars]
        file_item.metadata = meta
        file_item.status = FileStatus.EXTRACTED

    except Exception as e:
        logger.error(f"提取失败 {path.name}: {e}")
        file_item.status = FileStatus.FAILED
        file_item.error_message = str(e)


# ── PDF ─────────────────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> tuple[str, dict]:
    """Extract text and metadata from PDF files."""
    from PyPDF2 import PdfReader

    reader = PdfReader(str(path))
    meta = dict(reader.metadata) if reader.metadata else {}

    text_parts = []
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)
        if len("".join(text_parts)) > 10000:
            break  # Enough content

    text = "\n".join(text_parts)

    # Extract useful metadata
    clean_meta = {
        "pages": len(reader.pages),
        "author": str(meta.get("/Author", "")),
        "title": str(meta.get("/Title", "")),
        "subject": str(meta.get("/Subject", "")),
        "creator": str(meta.get("/Creator", "")),
        "producer": str(meta.get("/Producer", "")),
    }

    logger.debug(f"PDF提取: {path.name} — {len(text)} chars, {len(reader.pages)} pages")
    return text, clean_meta


# ── DOCX ────────────────────────────────────────────────────────────────────

def _extract_docx(path: Path) -> tuple[str, dict]:
    """Extract text from Word documents."""
    from docx import Document

    doc = Document(str(path))

    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
        if len("\n".join(text_parts)) > 8000:
            break

    text = "\n".join(text_parts)

    # Extract document properties
    props = doc.core_properties
    meta = {
        "author": str(props.author or ""),
        "title": str(props.title or ""),
        "subject": str(props.subject or ""),
        "category": str(props.category or ""),
        "created": str(props.created or ""),
        "modified": str(props.modified or ""),
        "paragraphs": len(doc.paragraphs),
    }

    logger.debug(f"DOCX提取: {path.name} — {len(text)} chars")
    return text, meta


# ── PPTX ────────────────────────────────────────────────────────────────────

def _extract_pptx(path: Path) -> tuple[str, dict]:
    """Extract text from PowerPoint presentations."""
    from pptx import Presentation

    prs = Presentation(str(path))

    text_parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        slide_text.append(p.text)
        if slide_text:
            text_parts.append("\n".join(slide_text))

    text = "\n---\n".join(text_parts)

    meta = {
        "slides": len(prs.slides),
    }

    logger.debug(f"PPTX提取: {path.name} — {len(text)} chars, {len(prs.slides)} slides")
    return text, meta


# ── XLSX ────────────────────────────────────────────────────────────────────

def _extract_xlsx(path: Path) -> tuple[str, dict]:
    """Extract text from Excel spreadsheets."""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)

    text_parts = []
    for sheet_name in wb.sheetnames[:5]:  # Max 5 sheets
        ws = wb[sheet_name]
        text_parts.append(f"[Sheet: {sheet_name}]")
        row_count = 0
        for row in ws.iter_rows(max_row=200, values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                text_parts.append(row_text)
                row_count += 1
            if row_count >= 100:
                break

    text = "\n".join(text_parts)

    meta = {
        "sheets": wb.sheetnames,
        "sheet_count": len(wb.sheetnames),
    }

    wb.close()
    logger.debug(f"XLSX提取: {path.name} — {len(text)} chars")
    return text, meta


# ── Plain Text / Markdown / HTML ────────────────────────────────────────────

def _extract_text(path: Path) -> tuple[str, dict]:
    """Extract text from plain text files (TXT, MD, CSV, HTML, RTF)."""
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    text = ""

    for enc in encodings:
        try:
            with open(path, "r", encoding=enc, errors="replace") as f:
                text = f.read()
            break
        except Exception:
            continue

    meta = {
        "encoding": enc if text else "unknown",
        "line_count": text.count("\n") if text else 0,
    }

    logger.debug(f"文本提取: {path.name} — {len(text)} chars")
    return text, meta


# ── Image ───────────────────────────────────────────────────────────────────

def _extract_image(path: Path) -> tuple[str, dict]:
    """Extract metadata from image files (EXIF, etc.). No OCR in MVP."""
    from PIL import Image
    from PIL.ExifTags import TAGS

    img = Image.open(str(path))

    # Extract EXIF
    exif_data = {}
    try:
        exif = img._getexif()
        if exif:
            for tag_id, value in exif.items():
                tag_name = TAGS.get(tag_id, tag_id)
                if isinstance(value, bytes):
                    value = value.decode("utf-8", errors="replace")[:200]
                exif_data[str(tag_name)] = str(value)
    except Exception:
        pass

    # Build descriptive text from metadata
    meta_text_parts = []
    if "DateTimeOriginal" in exif_data:
        meta_text_parts.append(f"拍摄日期: {exif_data['DateTimeOriginal']}")
    if "Make" in exif_data:
        meta_text_parts.append(f"相机: {exif_data['Make']} {exif_data.get('Model', '')}")
    if "ImageDescription" in exif_data:
        meta_text_parts.append(f"描述: {exif_data['ImageDescription']}")

    text = "\n".join(meta_text_parts) if meta_text_parts else f"图片文件: {path.name}\n尺寸: {img.size[0]}x{img.size[1]}"

    meta = {
        "width": img.size[0],
        "height": img.size[1],
        "format": img.format or path.suffix.upper().lstrip("."),
        "mode": img.mode,
        **{f"exif_{k}": v for k, v in list(exif_data.items())[:20]},
    }

    logger.debug(f"图片提取: {path.name} — {img.size[0]}x{img.size[1]}")
    return text, meta


# ── Video ───────────────────────────────────────────────────────────────────

def _extract_video(path: Path) -> tuple[str, dict]:
    """Extract basic info from video files (no frame analysis in MVP)."""
    stat = path.stat()

    text = f"视频文件: {path.name}\n大小: {stat.st_size / (1024*1024):.1f} MB\n修改时间: {stat.st_mtime}"

    meta = {
        "size_mb": round(stat.st_size / (1024 * 1024), 1),
        "format": path.suffix.upper().lstrip("."),
    }

    logger.debug(f"视频提取: {path.name}")
    return text, meta
